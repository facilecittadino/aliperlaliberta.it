# -*- coding: utf-8 -*-
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = Path(__file__).resolve().parent
LOGO = ROOT / "logo.png"

W = 7.5
H = 13.333
NAVY = RGBColor(10, 46, 92)
BLUE = RGBColor(11, 94, 215)
CYAN = RGBColor(28, 181, 214)
GREEN = RGBColor(39, 170, 116)
INK = RGBColor(17, 34, 54)
WHITE = RGBColor(255, 255, 255)
MIST = RGBColor(239, 246, 252)
PINK = RGBColor(229, 54, 115)

VIDEOS = [
    {
        "slug": "01-chi-siamo",
        "title": "Chi siamo",
        "caption": "Orientamento chiaro, ascolto e servizi con CAF e Patronato partner.",
        "slides": [
            ("Hai bisogno di orientamento?", "Partiamo dalla tua situazione, con calma.", NAVY, CYAN),
            ("Siamo un'associazione.", "Ti ascoltiamo e ti aiutiamo a capire documenti, diritti e prossimi passi.", WHITE, BLUE),
            ("Collaboriamo con CAF e Patronato partner.", "Per offrirti servizi fiscali e previdenziali con un percorso più semplice.", MIST, GREEN),
            ("Inizia dalla tua area cliente.", "aliperlaliberta.it\n351 365 7045", BLUE, WHITE),
        ],
    },
    {
        "slug": "02-documenti-pratica",
        "title": "Documenti per la pratica",
        "caption": "Una lista semplice per arrivare preparati al primo contatto.",
        "slides": [
            ("Quali documenti servono?", "Ogni pratica è diversa. Verifichiamo insieme ciò che serve davvero.", INK, CYAN),
            ("Porta con te", "Documento di identità\nCodice fiscale\nDocumenti relativi alla richiesta", WHITE, BLUE),
            ("Prima dell'appuntamento", "Controlla la lista nell'area cliente o chiedi al nostro assistente.", MIST, PINK),
            ("Meno dubbi. Più chiarezza.", "Prenota su aliperlaliberta.it\n351 365 7045", GREEN, WHITE),
        ],
    },
    {
        "slug": "03-segui-richiesta",
        "title": "Segui la tua richiesta",
        "caption": "Dall'invio agli aggiornamenti, tutto nella tua area riservata.",
        "slides": [
            ("La richiesta non finisce dopo l'invio.", "Puoi seguirla senza perdere informazioni.", NAVY, CYAN),
            ("Dalla tua area cliente puoi", "Prenotare servizi\nVedere le tue richieste\nControllare lo stato", WHITE, BLUE),
            ("I tuoi dati restano nella tua area riservata.", "Ogni persona vede soltanto le proprie richieste.", MIST, GREEN),
            ("Accedi. Invia. Segui.", "aliperlaliberta.it\n351 365 7045", BLUE, WHITE),
        ],
    },
]


def add_transition(slide, duration_ms=3000):
    transition = parse_xml(
        f'<p:transition {nsdecls("p")} advClick="0" advTm="{duration_ms}"><p:fade/></p:transition>'
    )
    slide._element.insert_element_before(transition, "p:timing", "p:extLst")


def add_text(slide, text, x, y, w, h, size, color, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    for index, line in enumerate(text.split("\n")):
        if index:
            paragraph.add_line_break()
        run = paragraph.add_run()
        run.text = line
        run.font.name = "Aptos Display"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_brand(slide, dark):
    plate = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.5), Inches(1.0), Inches(1.0))
    plate.fill.solid()
    plate.fill.fore_color.rgb = WHITE
    plate.line.fill.background()
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(0.64), Inches(0.59), width=Inches(0.82), height=Inches(0.82))
    label_color = WHITE if dark else NAVY
    add_text(slide, "ALI PER LA LIBERTA", 1.8, 0.56, 5.05, 0.78, 16, label_color, True, PP_ALIGN.RIGHT)


def add_progress(slide, index, total, accent, dark):
    base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(12.56), Inches(6.3), Inches(0.05))
    base.fill.solid()
    base.fill.fore_color.rgb = RGBColor(255, 255, 255) if dark else RGBColor(201, 216, 232)
    base.line.fill.background()
    fill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(12.56), Inches(6.3 * index / total), Inches(0.05))
    fill.fill.solid()
    fill.fill.fore_color.rgb = accent
    fill.line.fill.background()


def create_video_deck(video):
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]
    total = len(video["slides"])

    for index, (headline, body, background, accent) in enumerate(video["slides"], start=1):
        slide = prs.slides.add_slide(blank)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = background
        dark = background in (NAVY, BLUE, INK, GREEN)
        title_color = WHITE if dark else INK
        body_color = RGBColor(226, 238, 250) if dark else RGBColor(61, 79, 99)

        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.45), Inches(0.12), Inches(1.15))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent
        accent_bar.line.fill.background()

        add_brand(slide, dark)
        add_text(slide, headline, 0.95, 2.25, 5.9, 3.2, 38, title_color, True)
        add_text(slide, body, 0.95, 6.0, 5.65, 3.25, 23, body_color, False)

        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(10.75), Inches(3.2), Inches(0.72))
        badge.fill.solid()
        badge.fill.fore_color.rgb = accent
        badge.line.fill.background()
        add_text(slide, video["title"].upper(), 1.12, 10.84, 2.86, 0.46, 12, WHITE if accent != WHITE else NAVY, True)
        add_progress(slide, index, total, accent, dark)
        add_transition(slide, 3000)

    output = OUT / f'{video["slug"]}.pptx'
    prs.save(output)
    return output


if __name__ == "__main__":
    created = [create_video_deck(video) for video in VIDEOS]
    print("\n".join(str(path) for path in created))
